"""
PDF to PPT conversion service
This module handles the conversion of PDF files to PowerPoint format using PyMuPDF and python-pptx.
"""

import os
import logging
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
from typing import Optional, List, Tuple

logger = logging.getLogger(__name__)

class PDFToPptConverter:
    """Handles conversion of PDF files to PowerPoint format."""
    
    def __init__(self):
        """Initialize the PDF to PPT converter."""
        pass
    
    def convert(
        self, 
        pdf_path: str, 
        output_path: Optional[str] = None,
        pages: Optional[List[int]] = None
    ) -> str:
        """
        Convert a PDF file to PowerPoint format.
        
        Args:
            pdf_path: Path to the PDF file
            output_path: Path for the output PPTX file (if None, uses the same filename with .pptx extension)
            pages: Specific pages to convert (if None, converts all pages)
            
        Returns:
            Path to the generated PPTX file
        """
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
            
        # If output path is not provided, use the same name with .pptx extension
        if output_path is None:
            output_path = os.path.splitext(pdf_path)[0] + ".pptx"
            
        try:
            # Create a new presentation
            prs = Presentation()
            
            # Add slide layouts
            slide_layout = prs.slide_layouts[5]  # blank layout
            
            # Open the PDF file
            pdf = fitz.open(pdf_path)
            
            # Determine which pages to convert
            page_indices = list(range(pdf.page_count))
            if pages is not None:
                # Filter to only the requested pages (convert from 1-based to 0-based indices)
                page_indices = [i-1 for i in pages if 0 <= i-1 < pdf.page_count]
            
            # Process each page
            for page_idx in page_indices:
                # Get the PDF page
                pdf_page = pdf.load_page(page_idx)
                
                # Extract text
                text = pdf_page.get_text()
                
                # Convert page to image
                pix = pdf_page.get_pixmap(alpha=False)
                
                # Save the image to a temporary file
                temp_img_path = os.path.join(tempfile.gettempdir(), f"page_{page_idx}.png")
                pix.save(temp_img_path)
                
                # Add a slide
                slide = prs.slides.add_slide(slide_layout)
                
                # Add the image to the slide
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)  # Adjust as needed
                height = Inches(6.5)  # Adjust as needed
                
                slide.shapes.add_picture(temp_img_path, left, top, width, height)
                
                # Clean up temporary image file
                os.remove(temp_img_path)
                
                # Add a text box with some of the text (optional)
                # We're only adding a small portion to avoid overwhelming the slide
                if text:
                    txbox = slide.shapes.add_textbox(left, Inches(7.2), width, Inches(0.5))
                    tf = txbox.text_frame
                    tf.text = f"Page {page_idx + 1} - PDF Conversion"
            
            # Close the PDF file
            pdf.close()
            
            # Save the presentation
            prs.save(output_path)
            
            logger.info(f"Successfully converted {pdf_path} to {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error converting PDF to PPT: {str(e)}")
            raise
    
    @staticmethod
    def parse_page_range(page_string: str) -> List[int]:
        """
        Parse a page range string (e.g., "1,3-5,7") into a list of page numbers.
        
        Args:
            page_string: String representation of page ranges
            
        Returns:
            List of page numbers
        """
        if not page_string:
            return None
            
        pages = []
        parts = page_string.split(',')
        
        for part in parts:
            part = part.strip()
            if '-' in part:
                # Handle page range (e.g., "3-5")
                start, end = map(int, part.split('-'))
                pages.extend(range(start, end + 1))
            else:
                # Handle single page
                pages.append(int(part))
                
        return pages

def convert_pdf_to_ppt(
    pdf_path: str,
    output_path: Optional[str] = None,
    page_string: Optional[str] = None
) -> str:
    """
    Convenience function to convert a PDF file to PowerPoint format.
    
    Args:
        pdf_path: Path to the PDF file
        output_path: Path for the output PPTX file
        page_string: String representation of page ranges (e.g., "1,3-5,7")
        
    Returns:
        Path to the generated PPTX file
    """
    # Create converter instance
    converter = PDFToPptConverter()
    
    # Parse page ranges if provided
    pages = None
    if page_string:
        pages = converter.parse_page_range(page_string)
        
    # Convert the PDF to PPT
    return converter.convert(pdf_path, output_path, pages)
